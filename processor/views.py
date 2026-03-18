import json
import os
import time
import docx
import pdfplumber
import requests
from dotenv import load_dotenv
from pptx import Presentation
from django.conf import settings
from django.http import JsonResponse
from django.shortcuts import render
from django.views.decorators.csrf import csrf_exempt
from .models import Subject, Topic, SubTopic

load_dotenv()


def boogie_woogie(prompt):
    # Loop through tokens and models to find an available one. If status 429 is hit, try the next token/model combo.
    tokens = [t for t in [os.getenv("HF_TOKEN_1"), os.getenv("HF_TOKEN_2")] if t]
    models = [
        "Qwen/Qwen2.5-Coder-32B-Instruct",
        "deepseek-ai/DeepSeek-V3",
        "deepseek-ai/DeepSeek-R1-Distill-Qwen-32B",
    ]

    for token in tokens:
        headers = {"Authorization": f"Bearer {token}"}
        for model in models:
            try:
                res = requests.post(
                    "https://router.huggingface.co/v1/chat/completions",
                    headers=headers,
                    json={
                        "model": model,
                        "messages": [{"role": "user", "content": prompt}],
                        "temperature": 0.0,
                    },
                    timeout=60,
                )
                if res.status_code == 200:
                    return res.json()
                if res.status_code == 429:
                    break
            except Exception:
                continue
    return None


def calculate_difficulty(text):
    word_count = len(text.split())
    # 1 point for every 300 words, max 5.
    score = (word_count // 300) + 1
    return min(score, 5)


@csrf_exempt
def get_metadata(request):
    # Checks file type and tells the frontend how many pages/slides to process.
    if "file" not in request.FILES:
        return JsonResponse({"error": "No file"}, status=400)

    file = request.FILES["file"]
    ext = file.name.lower()

    try:
        if ext.endswith(".pdf"):
            with pdfplumber.open(file) as pdf:
                count, unit = len(pdf.pages), "pages"
        elif ext.endswith(".docx"):
            doc = docx.Document(file)
            count, unit = sum(1 for p in doc.paragraphs if p.text.strip()), "paragraphs"
        elif ext.endswith(".pptx"):
            prs = Presentation(file)
            count, unit = len(prs.slides), "slides"
        else:
            return JsonResponse({"error": "Unsupported file"}, status=400)

        return JsonResponse({"total": count, "type": unit})
    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)


@csrf_exempt
def home(request):
    return render(request, "index.html")


@csrf_exempt
def upload_content(request):
    if request.method != "POST":
        return JsonResponse({"error": "POST only"}, status=405)

    # Setup Context & Range
    context_subject = request.POST.get("document_context", "").strip()
    try:
        start_pg = max(0, int(request.POST.get("start_page", 1)) - 1)
        end_pg = int(request.POST.get("end_page", 3))
    except ValueError:
        start_pg, end_pg = 0, 3

    # Extract Text
    text = ""
    file = request.FILES.get("file")

    if file:
        ext = file.name.lower()
        try:
            if ext.endswith(".pdf"):
                with pdfplumber.open(file) as pdf:
                    pages = pdf.pages[start_pg : min(end_pg, len(pdf.pages))]
                    text = "\n".join(
                        [p.extract_text() for p in pages if p.extract_text()]
                    )
            elif ext.endswith(".docx"):
                doc = docx.Document(file)
                text = "\n".join([p.text for p in doc.paragraphs[start_pg:end_pg]])
            elif ext.endswith(".pptx"):
                prs = Presentation(file)
                slides = list(prs.slides)[start_pg:end_pg]
                text = "\n".join(
                    [
                        shape.text
                        for s in slides
                        for shape in s.shapes
                        if hasattr(shape, "text")
                    ]
                )
        except Exception as e:
            return JsonResponse({"error": f"File Error: {str(e)}"}, status=500)
    else:
        text = request.POST.get("manual_text", "")

    text = text.replace("\x00", "").strip()[:12000]
    if not text:
        return JsonResponse({"error": "No text found"}, status=400)

    # Prompt
    context_instruction = (
        f"IMPORTANT: This is a continuation of the document: '{context_subject}'. Keep the 'subject' field as '{context_subject}'."
        if context_subject
        else "Identify the overall Subject of this text."
    )

    prompt = f"""
    Analyze this PARTIAL text from a document. {context_instruction}
    Create a nested study guide: Subject > Topic > Subtopic > Content.
    CRITICAL: DO NOT add external info, code, or examples. Use ONLY the provided text.

    HIERARCHY LOGIC:
    1. 'Subject' is the overall title of the document.
    2. 'Topics' are the major sections (usually numbered or in all-caps). If a topic spans multiple batches, use the SAME 'topic_name'.
    3. 'Subtopics' are the specific concepts discussed within a topic.
    4. 'Content' is the explanation or definition for that specific subtopic.
    5. IMPORTANT: If a list of items has NO individual descriptions in the text, DO NOT make them separate subtopics. Instead, group them into a single Subtopic called "Key Concepts" or "Overview".

    STRICT JSON OUTPUT:
    {{
    "subject": "String",
    "topics": [
        {{
        "topic_name": "String",
        "subtopics": [
            {{ "subtopic_name": "String", "content": "String" }}
        ]
        }}
    ]
    }}

    TEXT TO PROCESS:
    {text}
    """

    ai_response_data = boogie_woogie(prompt)
    if not ai_response_data:
        return JsonResponse({"error": "AI Busy"}, status=503)

    # Save to Database
    try:
        raw_content = ai_response_data["choices"][0]["message"]["content"]
        start_idx, end_idx = raw_content.find("{"), raw_content.rfind("}")
        structured_data = json.loads(raw_content[start_idx : end_idx + 1])

        ai_name = structured_data.get("subject", "General Studies")
        final_name = (
            context_subject
            if context_subject and context_subject.lower() != "null"
            else ai_name
        )

        subj_obj, _ = Subject.objects.get_or_create(name=final_name)

        last_topic = None
        for t_data in structured_data.get("topics", []):
            topic_obj, _ = Topic.objects.get_or_create(
                subject=subj_obj, name=t_data.get("topic_name", "Untitled")
            )
            last_topic = topic_obj

            for s_data in t_data.get("subtopics", []):
                s_name = s_data.get("subtopic_name", "General Concept")
                SubTopic.objects.get_or_create(
                    topic=topic_obj,
                    name=s_name,
                    defaults={
                        "content": s_data.get("content", ""),
                        "weight": calculate_difficulty(s_data.get("content", "")),
                    },
                )

        return JsonResponse(
            {
                "subject": subj_obj.name,
                "topic_name": last_topic.name if last_topic else "Uploaded",
                "subtopics": (
                    list(last_topic.subtopics.all().values("name", "content"))
                    if last_topic
                    else []
                ),
            }
        )

    except Exception as e:
        return JsonResponse({"error": f"Database Error: {str(e)}"}, status=500)


@csrf_exempt
def generate_quiz(request):
    # Generates 5 questions based on the content of a topic.
    topic_id = request.POST.get("guide_id")
    topic = Topic.objects.filter(id=topic_id).first()
    if not topic:
        return JsonResponse({"error": "Topic not found"}, status=404)

    lesson_context = "\n".join(
        [f"Subtopic: {s.name}\nContent: {s.content}" for s in topic.subtopics.all()]
    )

    quiz_prompt = f"""
        Based on this lesson content, generate a 5-question Multiple Choice Quiz.
        LESSON CONTENT: {lesson_context}
        
        STRICT JSON OUTPUT FORMAT:
        {{
            "quiz_title": "{topic.name} Quiz",
            "questions": [
                {{
                    "question": "String",
                    "options": ["A", "B", "C", "D"],
                    "correct_answer": "Exact string from options"
                }}
            ]
        }}
        """

    ai_data = boogie_woogie(quiz_prompt)
    if not ai_data:
        return JsonResponse({"error": "AI busy"}, status=503)

    try:
        raw = ai_data["choices"][0]["message"]["content"]
        return JsonResponse(json.loads(raw[raw.find("{") : raw.rfind("}") + 1]))
    except Exception:
        return JsonResponse({"error": "AI Parse Error"}, status=500)


@csrf_exempt
def get_all_lessons(request):
    # Prefetch related to minimize DB hits. Difficulty is averaged across subtopics for a topic-level score.
    subjects = Subject.objects.prefetch_related("topics__subtopics").all()

    data = []
    for s in subjects:
        topics = []
        for t in s.topics.all():
            # Calculate average difficulty for the topic based on its subtopics
            weights = [st.weight for st in t.subtopics.all()]
            avg = sum(weights) // len(weights) if weights else 1
            topics.append({"id": t.id, "name": t.name, "difficulty": avg})

        data.append({"subject_name": s.name, "topics": topics})

    return JsonResponse(data, safe=False)


@csrf_exempt
def delete_lesson(request):
    topic = Topic.objects.filter(id=request.POST.get("guide_id")).first()
    if not topic:
        return JsonResponse({"error": "Not found"}, status=404)

    subject = topic.subject
    topic.delete()

    # If the subject has no more topics, delete it as well to keep the UI clean.
    if not subject.topics.exists():
        subject.delete()

    return JsonResponse({"success": True})


@csrf_exempt
def get_lesson_detail(request):
    # Returns the full content of a lesson for the detail view.
    topic = Topic.objects.filter(id=request.POST.get("guide_id")).first()
    if not topic:
        return JsonResponse({"error": "Not found"}, status=404)

    return JsonResponse(
        {
            "topic_name": topic.name,
            "subtopics": list(
                topic.subtopics.all().values("name", "content", "weight")
            ),
        }
    )
