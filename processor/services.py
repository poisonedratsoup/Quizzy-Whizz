import json
import os
import io
import docx
import pdfplumber
import requests
from pptx import Presentation
from dotenv import load_dotenv

load_dotenv()


def boogie_woogie(prompt):
    # Try multiple tokens/models to handle rate limits
    tokens = [t for t in [os.getenv("HF_TOKEN_1"), os.getenv("HF_TOKEN_2")] if t]
    models = [
        "Qwen/Qwen2.5-Coder-32B-Instruct",
        "deepseek-ai/DeepSeek-V3",
        "deepseek-ai/DeepSeek-R1-Distill-Qwen-32B",
    ]
    
    if not tokens:
        print("Error: No Hugging Face tokens found in environment.")
        return None
    
    for token in tokens:
        headers = {"Authorization": f"Bearer {token}"}
        for model in models:
            try:
                res = requests.post(
                    "https://router.huggingface.co/v1/chat/completions",
                    headers=headers,
                    json={
                        "model": model,
                        "messages": [{"role": "user", "content": prompt}],
                        "temperature": 0.0,
                    },
                    timeout=60,
                )
                if res.status_code == 200:
                    return res.json()
                if res.status_code == 429:
                    print(f"Token rate limited for model {model}. Switching token...")
                    break
                
                print(f"Model {model} failed with status {res.status_code}. Trying next model...")
            except Exception as e:
                print(f"Connection error with {model}: {e}")
                continue
    return None


def extract_text_from_file(file):
    # Only process first few pages/paragraphs to avoid token overload
    ext = file.name.lower()
    text = ""
    file_content = io.BytesIO(file.read())

    try:
        if ext.endswith(".pdf"):
            with pdfplumber.open(file_content) as pdf:
                text = "\n".join(
                    [p.extract_text() for p in pdf.pages[:3] if p.extract_text()]
                )
        elif ext.endswith(".docx"):
            doc = docx.Document(file_content)
            text = "\n".join([p.text for p in doc.paragraphs[:50]])
        elif ext.endswith(".pptx"):
            prs = Presentation(file_content)
            text = "\n".join(
                [
                    shape.text
                    for s in list(prs.slides)[:3]
                    for shape in s.shapes
                    if hasattr(shape, "text")
                ]
            )
        elif ext.endswith(".txt"):
            text = file_content.getvalue().decode("utf-8")
    except Exception as e:
        print(f"File extraction failed: {e}")
        
    return text.replace("\x00", "").strip()[:10000]


def calculate_difficulty(text):
    # Calculate difficulty based on word count, capped at 5
    word_count = len(text.split())
    return min((word_count // 200) + 1, 5)


def parse_ai_json(ai_res):
    # Extract JSON from AI response, handling potential formatting issues
    try:
        raw = ai_res["choices"][0]["message"]["content"]
        start_idx, end_idx = raw.find("{"), raw.rfind("}")
        if start_idx != -1 and end_idx != -1:
            return json.loads(raw[start_idx : end_idx + 1])
    except (KeyError, ValueError, TypeError) as e:
        print(f"JSON Parsing failed: {e}")
    return None


def normalize_name(name):
    # Remove dashes, extra spaces, and make lowercase
    return name.lower().replace("-", " ").replace("_", " ").strip()
