from flask import Flask, request, jsonify
from flask_cors import CORS
import pdfplumber
import requests
import json
import os
from dotenv import load_dotenv
import docx
from pptx import Presentation

load_dotenv()

HF_TOKEN = os.getenv("HF_TOKEN")

app = Flask(__name__)
CORS(app) 

HF_MODEL_URL = "https://router.huggingface.co/v1/chat/completions"

@app.route('/get_pages', methods=['POST'])
def get_pages():
    if 'file' not in request.files:
        return jsonify({"error": "No file"}), 400
    file = request.files['file']
    try:
        with pdfplumber.open(file) as pdf:
            return jsonify({"total_pages": len(pdf.pages)})
    except Exception as e:
        return jsonify({"error": "Could not read PDF"}), 500

@app.route('/upload', methods=['POST'])
def upload_pdf():
    text = ""
    filename = ""

    try:
        start_val = int(request.form.get('start_page', 1))
        end_val = int(request.form.get('end_page', 3))
        start_pg = max(0, start_val - 1)
        end_pg = end_val
    except ValueError:
        start_pg, end_pg = 0, 3
    
    if 'file' in request.files:
        file = request.files['file']
        if file.filename == '':
             return jsonify({"error": "No file selected"}), 400

        filename = file.filename.lower()
        if filename.endswith('.pdf'):
            with pdfplumber.open(file) as pdf:
                total = len(pdf.pages)
                actual_end = min(end_pg, total)
                text = "\n".join([page.extract_text() for page in pdf.pages[start_pg:actual_end] if page.extract_text()])

        elif filename.endswith('.docx'):
            doc = docx.Document(file)
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            paras_per_page = 15
            start_idx = start_pg * paras_per_page
            end_idx = end_pg * paras_per_page
            text = "\n".join(paragraphs[start_idx:end_idx])

        elif filename.endswith('.pptx'):
            prs = Presentation(file)
            all_slides = list(prs.slides) 
            selected_slides = all_slides[start_pg:end_pg]
            
            full_text = []
            for slide in selected_slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text.append(shape.text)
            text = "\n".join(full_text)

    elif 'manual_text' in request.form:
        text = request.form['manual_text'] 
    
    text = text.replace('\x00', '').strip()[:12000]
    if not text.strip():
        return jsonify({
            "error": f"No text found in Batch {start_val}-{end_val}. The slides might be empty or contains only images."
        }), 400
    
    strict_rule = ""
    if filename.endswith(('.pptx', '.docx')):
        strict_rule = "CRITICAL: DO NOT add external info, code, or examples. Use ONLY the provided text."

    prompt = f"""
    {strict_rule}
    Analyze the following text and extract its hierarchical structure.
    Create a nested study guide: Subject > Topic > Subtopic > Content.

    HIERARCHY LOGIC:
    1. 'Subject' is the overall title of the document.
    2. 'Topics' are the major sections (usually numbered or in all-caps).
    3. 'Subtopics' are the specific concepts discussed within a topic.
    4. 'Content' is the explanation or definition for that specific subtopic.
    5. IMPORTANT: If a list of items has NO individual descriptions in the text, DO NOT make them separate subtopics. Instead, group them into a single Subtopic called "Key Concepts" or "Overview".

    STRICT JSON OUTPUT:
    {{
    "subject": "String",
    "topics": [
        {{
        "topic_name": "String",
        "subtopics": [
            {{ "subtopic_name": "String", "content": "String" }}
        ]
        }}
    ]
    }}

    TEXT TO PROCESS:
    {text}
    """

    payload = {
        "model": "Qwen/Qwen2.5-Coder-32B-Instruct", 
        "messages": [
            {"role": "system", "content": "You are a helpful tutor that outputs ONLY valid JSON."},
            {"role": "user", "content": prompt}
        ],
        "temperature": 0.0
    }
    
    response = requests.post(HF_MODEL_URL, headers={"Authorization": f"Bearer {HF_TOKEN}"}, json=payload, timeout=(5, 60))
    if response.status_code != 200:
        print(f" HF API Error: {response.status_code} - {response.text}")
        return jsonify({"error": "The AI is currently overwhelmed. Please wait a few seconds and try again!"}), 503
    
    result = response.json()
    ai_text = result["choices"][0]["message"]["content"]
    
    start_idx = ai_text.find('{')
    end_idx = ai_text.rfind('}')

    if start_idx == -1 or end_idx == -1:
        return jsonify({"error": "AI response was not in the correct format."}), 500

    json_string = ai_text[start_idx : end_idx + 1]
    
    try:
        return json.loads(json_string)
    except json.JSONDecodeError:
        return jsonify({"error": "The AI created a broken structure. Please try again."}), 500

if __name__ == '__main__':
    app.run(debug=True, port=5000)
